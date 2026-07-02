"""Document parsers (§5): PDF / DOCX / PPTX / XLSX → normalized text + tables.

Lightweight, dependency-safe parsers (pypdf/pdfplumber, python-docx, python-pptx,
openpyxl) behind one ``parse_document`` entry point. Robust to corrupt files
(returns ``None`` and logs) per §24.17.
"""

from __future__ import annotations

import hashlib
import re
from dataclasses import dataclass, field
from pathlib import Path

from kg_common import get_logger

_log = get_logger("parsers")

SUPPORTED = {".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".txt", ".md"}


@dataclass
class ParsedTable:
    page: int
    rows: list[list[str]]


@dataclass
class ParsedDoc:
    path: str
    title: str
    doc_type: str
    file_hash: str
    lang: str
    pages: list[tuple[int, str]] = field(default_factory=list)  # (page_no, text)
    tables: list[ParsedTable] = field(default_factory=list)
    country: str | None = None
    year: int | None = None

    @property
    def full_text(self) -> str:
        return "\n\n".join(t for _, t in self.pages)


_YEAR_RE = re.compile(r"\b(19[7-9]\d|20[0-4]\d)\b")


def _file_hash(path: Path) -> str:
    h = hashlib.sha1()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1 << 16), b""):
            h.update(chunk)
    return h.hexdigest()[:16]


def _detect_lang(text: str) -> str:
    cyr = len(re.findall(r"[а-яё]", text, re.I))
    lat = len(re.findall(r"[a-z]", text, re.I))
    if not (cyr or lat):
        return "unknown"
    if cyr and lat and min(cyr, lat) / max(cyr, lat) > 0.3:
        return "mixed"
    return "ru" if cyr >= lat else "en"


def _guess_doc_type(name: str, parent: str) -> str:
    p = parent.lower()
    if "патент" in p or "patent" in p:
        return "patent"
    if "обзор" in p or "review" in p:
        return "review"
    if "статьи" in p or "article" in p or "журнал" in p:
        return "article"
    if "доклад" in p or "презент" in p or ".pptx" in name.lower():
        return "presentation"
    if "конференц" in p or "conference" in p:
        return "conference"
    if "отчет" in p or "отчёт" in p or "report" in p:
        return "internal_report"
    return "article"


def _guess_country(text: str) -> str | None:
    from kg_schema.taxonomy import load_taxonomy

    idx = load_taxonomy()
    low = text[:4000].lower()
    for e in idx.entries:
        if e.node_type == "Country":
            for term in (e.canonical_ru, e.canonical_en, *e.aliases):
                if term and term.lower() in low:
                    return e.id
    return None


def parse_document(path: str | Path) -> ParsedDoc | None:
    p = Path(path)
    ext = p.suffix.lower()
    if ext not in SUPPORTED:
        return None
    try:
        if ext == ".pdf":
            pages, tables = _parse_pdf(p)
        elif ext == ".docx":
            pages, tables = _parse_docx(p)
        elif ext == ".pptx":
            pages, tables = _parse_pptx(p)
        elif ext in (".xlsx", ".xls"):
            pages, tables = _parse_xlsx(p)
        else:
            pages, tables = [(1, p.read_text(encoding="utf-8", errors="ignore"))], []
    except Exception as exc:
        _log.warning("parse.failed", path=str(p), error=str(exc)[:150])
        return None

    text = "\n".join(t for _, t in pages)
    if len(text.strip()) < 30:
        return None
    ym = _YEAR_RE.search(p.name) or _YEAR_RE.search(text[:2000])
    return ParsedDoc(
        path=str(p),
        title=p.stem,
        doc_type=_guess_doc_type(p.name, p.parent.name),
        file_hash=_file_hash(p),
        lang=_detect_lang(text),
        pages=pages,
        tables=tables,
        country=_guess_country(text),
        year=int(ym.group(0)) if ym else None,
    )


def _parse_pdf(p: Path) -> tuple[list[tuple[int, str]], list[ParsedTable]]:
    pages: list[tuple[int, str]] = []
    tables: list[ParsedTable] = []
    try:
        import pdfplumber

        with pdfplumber.open(str(p)) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                txt = page.extract_text() or ""
                if txt.strip():
                    pages.append((i, txt))
                for tbl in page.extract_tables() or []:
                    rows = [[(c or "").strip() for c in row] for row in tbl if row]
                    if rows:
                        tables.append(ParsedTable(page=i, rows=rows))
                if i >= 60:  # cap very large PDFs
                    break
    except Exception:
        from pypdf import PdfReader

        reader = PdfReader(str(p))
        for i, page in enumerate(reader.pages[:60], start=1):
            txt = page.extract_text() or ""
            if txt.strip():
                pages.append((i, txt))
    return pages, tables


def _parse_docx(p: Path) -> tuple[list[tuple[int, str]], list[ParsedTable]]:
    import docx

    d = docx.Document(str(p))
    paras = [para.text for para in d.paragraphs if para.text.strip()]
    tables = []
    for t in d.tables:
        rows = [[c.text.strip() for c in row.cells] for row in t.rows]
        if rows:
            tables.append(ParsedTable(page=1, rows=rows))
    text = "\n".join(paras)
    # split into pseudo-pages of ~3000 chars
    pages = [(i + 1, text[j : j + 3000]) for i, j in enumerate(range(0, len(text), 3000))]
    return pages or [(1, text)], tables


def _parse_pptx(p: Path) -> tuple[list[tuple[int, str]], list[ParsedTable]]:
    from pptx import Presentation

    prs = Presentation(str(p))
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        txt = "\n".join(x for x in parts if x.strip())
        if txt.strip():
            pages.append((i, txt))
    return pages, []


def _parse_xlsx(p: Path) -> tuple[list[tuple[int, str]], list[ParsedTable]]:
    import openpyxl

    wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
    pages = []
    tables = []
    for i, ws in enumerate(wb.worksheets, start=1):
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(cells)
            if len(rows) >= 200:
                break
        if rows:
            tables.append(ParsedTable(page=i, rows=rows))
            pages.append((i, "\n".join(" | ".join(r) for r in rows)))
    wb.close()
    return pages, tables
